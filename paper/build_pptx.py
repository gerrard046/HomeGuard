#!/usr/bin/env python3
"""Generator slide presentasi (PowerPoint) untuk paparan paper HomeGuard.

Menghasilkan ``paper/HomeGuard_paparan.pptx`` — siap dipakai saat bimbingan,
seminar, atau sidang. Memakai pustaka ``python-pptx``:

    pip install python-pptx
    python paper/build_pptx.py        # -> paper/HomeGuard_paparan.pptx

Slide disusun mengikuti alur: judul → latar belakang → tujuan → arsitektur
(pipeline 4 modul) → cara kerja → evolusi prototipe (M1 vs M12) → fitur →
justifikasi vs alat lain → pengujian → demo → kesimpulan → penutup.
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

PAPER_DIR = os.path.dirname(os.path.abspath(__file__))
TUJUAN = os.path.join(PAPER_DIR, "HomeGuard_paparan.pptx")

# Palet warna (tema keamanan: biru tua + aksen hijau & merah)
BIRU_TUA = RGBColor(0x0D, 0x1B, 0x3E)
BIRU = RGBColor(0x1B, 0x3A, 0x6B)
PUTIH = RGBColor(0xFF, 0xFF, 0xFF)
ABU = RGBColor(0x33, 0x3A, 0x45)
ABU_TERANG = RGBColor(0x6B, 0x72, 0x80)
HIJAU = RGBColor(0x1A, 0x9E, 0x5A)
MERAH = RGBColor(0xD1, 0x33, 0x33)
ORANYE = RGBColor(0xE0, 0x8A, 0x1E)
AKSEN = RGBColor(0x2E, 0x86, 0xDE)

LEBAR = Inches(13.333)
TINGGI = Inches(7.5)


def _set_bg(slide, warna):
    """Warnai latar belakang penuh sebuah slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = warna


def _kotak(slide, l, t, w, h, warna, garis=None):
    """Gambar persegi panjang berisi warna solid; kembalikan shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = warna
    if garis is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = garis
        shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp


def _teks(slide, l, t, w, h, teks, ukuran=18, warna=ABU, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
          font="Calibri"):
    """Tambah text box dengan satu paragraf bergaya."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = teks
    r.font.size = Pt(ukuran)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = warna
    r.font.name = font
    return tb


def _poin(slide, l, t, w, h, item, ukuran=18, warna=ABU, gap=10,
          warna_bullet=AKSEN):
    """Tambah daftar berpoin. ``item`` = list[str] atau list[(teks, level)]."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(item):
        if isinstance(it, tuple):
            txt, level = it
        else:
            txt, level = it, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = level
        bullet = "•  " if level == 0 else "–  "
        rb = p.add_run()
        rb.text = bullet
        rb.font.size = Pt(ukuran)
        rb.font.bold = True
        rb.font.color.rgb = warna_bullet
        rb.font.name = "Calibri"
        rt = p.add_run()
        rt.text = txt
        rt.font.size = Pt(ukuran if level == 0 else ukuran - 2)
        rt.font.color.rgb = warna
        rt.font.name = "Calibri"
    return tb


def _header(slide, judul, nomor):
    """Bar judul atas + nomor slide untuk slide konten."""
    _kotak(slide, 0, 0, LEBAR, Inches(1.05), BIRU_TUA)
    _kotak(slide, 0, Inches(1.05), LEBAR, Inches(0.06), HIJAU)
    _teks(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
          judul, ukuran=28, warna=PUTIH, bold=True,
          anchor=MSO_ANCHOR.MIDDLE)
    _teks(slide, Inches(12.4), Inches(0.18), Inches(0.7), Inches(0.7),
          str(nomor), ukuran=14, warna=RGBColor(0x9F, 0xB3, 0xD1),
          bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _slide_kosong(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# --------------------------------------------------------------------------- #
#  Pembangun tiap slide
# --------------------------------------------------------------------------- #

def slide_judul(prs):
    s = _slide_kosong(prs)
    _set_bg(s, BIRU_TUA)
    _kotak(s, 0, Inches(2.7), LEBAR, Inches(0.08), HIJAU)
    _teks(s, Inches(1.0), Inches(0.9), Inches(11.3), Inches(0.5),
          "POLITEKNIK SIBER DAN SANDI NEGARA", ukuran=18,
          warna=RGBColor(0x9F, 0xB3, 0xD1), bold=True, align=PP_ALIGN.CENTER)
    _teks(s, Inches(1.0), Inches(2.85), Inches(11.3), Inches(0.6),
          "HomeGuard", ukuran=54, warna=PUTIH, bold=True,
          align=PP_ALIGN.CENTER)
    _teks(s, Inches(1.0), Inches(1.75), Inches(11.3), Inches(1.0),
          "Pengembangan Aplikasi Pemindai Keamanan Jaringan Rumah untuk "
          "Identifikasi Kerentanan Perangkat IoT pada Jaringan Wi-Fi Domestik",
          ukuran=20, warna=RGBColor(0xCF, 0xDD, 0xF0), bold=False,
          align=PP_ALIGN.CENTER)
    _teks(s, Inches(1.0), Inches(3.95), Inches(11.3), Inches(0.5),
          "Prototipe pemindai keamanan IoT berbasis OWASP IoT Top 10",
          ukuran=18, warna=HIJAU, italic=True, align=PP_ALIGN.CENTER)
    _teks(s, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.5),
          "Reiza Gerrard Rizki Ramadhan", ukuran=20, warna=PUTIH,
          bold=True, align=PP_ALIGN.CENTER)
    _teks(s, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.5),
          "Paparan Tugas Akhir", ukuran=16,
          warna=RGBColor(0x9F, 0xB3, 0xD1), align=PP_ALIGN.CENTER)


def slide_agenda(prs):
    s = _slide_kosong(prs)
    _set_bg(s, PUTIH)
    _header(s, "Agenda", 2)
    kiri = [
        "Latar Belakang & Permasalahan",
        "Tujuan Penelitian",
        "Arsitektur Aplikasi (Pipeline 4 Modul)",
        "Cara Kerja Sistem",
    ]
    kanan = [
        "Evolusi Prototipe (Awal → Saat Ini)",
        "Justifikasi & Diferensiasi",
        "Pengujian & Hasil",
        "Demo, Kesimpulan & Roadmap",
    ]
    for i, t in enumerate(kiri):
        y = Inches(1.6 + i * 1.15)
        _kotak(s, Inches(0.8), y, Inches(0.55), Inches(0.55), AKSEN)
        _teks(s, Inches(0.8), y, Inches(0.55), Inches(0.55), str(i + 1),
              ukuran=20, warna=PUTIH, bold=True, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
        _teks(s, Inches(1.55), y, Inches(5.0), Inches(0.55), t, ukuran=18,
              warna=ABU, anchor=MSO_ANCHOR.MIDDLE)
    for i, t in enumerate(kanan):
        y = Inches(1.6 + i * 1.15)
        _kotak(s, Inches(7.0), y, Inches(0.55), Inches(0.55), HIJAU)
        _teks(s, Inches(7.0), y, Inches(0.55), Inches(0.55), str(i + 5),
              ukuran=20, warna=PUTIH, bold=True, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
        _teks(s, Inches(7.75), y, Inches(5.0), Inches(0.55), t, ukuran=18,
              warna=ABU, anchor=MSO_ANCHOR.MIDDLE)


def slide_latar(prs):
    s = _slide_kosong(prs)
    _set_bg(s, PUTIH)
    _header(s, "Latar Belakang & Permasalahan", 3)
    _poin(s, Inches(0.8), Inches(1.4), Inches(7.0), Inches(5.5), [
        "Perangkat IoT rumahan (kamera CCTV, router, smart TV) tumbuh pesat",
        "Sering dikirim dengan konfigurasi default tidak aman:",
        ("port manajemen terbuka (Telnet, UPnP)", 1),
        ("kredensial bawaan admin/admin", 1),
        ("protokol & enkripsi lemah", 1),
        "Dieksploitasi botnet seperti Mirai (port 23/2323)",
        "Pengguna awam tidak punya alat yang mudah, lokal, & edukatif "
        "untuk memeriksa keamanannya sendiri",
    ], ukuran=18, gap=12)
    # Kartu statistik ancaman
    _kotak(s, Inches(8.2), Inches(1.6), Inches(4.3), Inches(2.3), BIRU_TUA)
    _teks(s, Inches(8.2), Inches(1.85), Inches(4.3), Inches(0.6), "Mirai",
          ukuran=24, warna=MERAH, bold=True, align=PP_ALIGN.CENTER)
    _teks(s, Inches(8.4), Inches(2.5), Inches(3.9), Inches(1.3),
          "Membajak ratusan ribu perangkat IoT lewat kredensial default — "
          "studi kasus utama pada paper ini",
          ukuran=15, warna=PUTIH, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)
    _kotak(s, Inches(8.2), Inches(4.1), Inches(4.3), Inches(2.3), BIRU)
    _teks(s, Inches(8.4), Inches(4.35), Inches(3.9), Inches(2.0),
          "Celah: alat keamanan IoT yang ada umumnya berbayar, berbasis "
          "cloud, atau terlalu teknis untuk pengguna rumahan.",
          ukuran=16, warna=PUTIH, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)


def slide_tujuan(prs):
    s = _slide_kosong(prs)
    _set_bg(s, PUTIH)
    _header(s, "Tujuan Penelitian", 4)
    kartu = [
        ("Membangun prototipe", "Pemindai keamanan jaringan rumah berbasis "
         "Python dengan pustaka standar (tanpa root, tanpa dependensi wajib)",
         AKSEN),
        ("Memetakan ke OWASP", "Mengaitkan temuan ke OWASP IoT Top 10 (2018) "
         "dan menghitung skor risiko 0–100 yang mudah dipahami", HIJAU),
        ("Edukatif & lokal", "Memberi pengguna awam laporan jelas berbahasa "
         "Indonesia, berjalan lokal demi privasi (tanpa cloud)", ORANYE),
    ]
    for i, (judul, isi, warna) in enumerate(kartu):
        x = Inches(0.8 + i * 4.05)
        _kotak(s, x, Inches(1.7), Inches(3.75), Inches(4.5),
               RGBColor(0xF2, 0xF5, 0xFA), garis=warna)
        _kotak(s, x, Inches(1.7), Inches(3.75), Inches(0.18), warna)
        _teks(s, x, Inches(2.1), Inches(3.75), Inches(0.7), judul, ukuran=21,
              warna=warna, bold=True, align=PP_ALIGN.CENTER)
        _teks(s, Inches(x.inches + 0.25), Inches(3.0), Inches(3.25),
              Inches(3.0), isi, ukuran=17, warna=ABU, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.TOP)


def slide_arsitektur(prs):
    s = _slide_kosong(prs)
    _set_bg(s, PUTIH)
    _header(s, "Arsitektur: Pipeline 4 Modul", 5)
    _teks(s, Inches(0.8), Inches(1.25), Inches(11.7), Inches(0.5),
          "Empat modul berurutan — keluaran satu modul menjadi masukan "
          "berikutnya (seperti ban berjalan):", ukuran=17, warna=ABU_TERANG)
    img = os.path.join(PAPER_DIR, "diagram_pipeline.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(0.9), Inches(2.0), width=Inches(11.5))
    _teks(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8),
          "Inti pemindaian hanya memakai pustaka standar Python — modul "
          "lanjutan (UDP, kredensial, TLS) bersifat opt-in.",
          ukuran=15, warna=ABU_TERANG, italic=True, align=PP_ALIGN.CENTER)


def slide_caramodul(prs):
    s = _slide_kosong(prs)
    _set_bg(s, PUTIH)
    _header(s, "Cara Kerja Tiap Modul", 6)
    mod = [
        ("1", "Network Discovery", AKSEN,
         "TCP ping (tanpa root) • parsing tabel ARP • lookup vendor dari "
         "MAC/OUI • ekspansi subnet CIDR"),
        ("2", "Port Scanning", HIJAU,
         "TCP connect paralel (ThreadPool) • probe UDP SSDP & mDNS • "
         "Nmap opsional dengan fallback socket"),
        ("3", "Service Detection", ORANYE,
         "Banner grabbing HTTP/Telnet/FTP/RTSP • deteksi versi • "
         "klasifikasi tipe perangkat (kamera/router/printer/NAS)"),
        ("4", "OWASP Mapping + Skor", MERAH,
         "Petakan temuan ke I1–I10 • cek kredensial bawaan, TLS, header • "
         "skor risiko = min(100, B_max + 3×n)"),
    ]
    for i, (no, judul, warna, isi) in enumerate(mod):
        y = Inches(1.45 + i * 1.4)
        _kotak(s, Inches(0.8), y, Inches(0.95), Inches(1.15), warna)
        _teks(s, Inches(0.8), y, Inches(0.95), Inches(1.15), no, ukuran=36,
              warna=PUTIH, bold=True, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
        _kotak(s, Inches(1.95), y, Inches(10.55), Inches(1.15),
               RGBColor(0xF2, 0xF5, 0xFA))
        _teks(s, Inches(2.2), Inches(y.inches + 0.12), Inches(10.0),
              Inches(0.5), judul, ukuran=20, warna=warna, bold=True)
        _teks(s, Inches(2.2), Inches(y.inches + 0.6), Inches(10.1),
              Inches(0.5), isi, ukuran=15, warna=ABU)


def slide_rumus(prs):
    s = _slide_kosong(prs)
    _set_bg(s, BIRU_TUA)
    _kotak(s, 0, 0, LEBAR, Inches(1.05), BIRU)
    _kotak(s, 0, Inches(1.05), LEBAR, Inches(0.06), HIJAU)
    _teks(s, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
          "Skor Risiko & Tingkat Severity", ukuran=28, warna=PUTIH,
          bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _teks(s, Inches(12.4), Inches(0.18), Inches(0.7), Inches(0.7), "7",
          ukuran=14, warna=RGBColor(0x9F, 0xB3, 0xD1), bold=True,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    _kotak(s, Inches(2.4), Inches(1.7), Inches(8.5), Inches(1.4), BIRU)
    _teks(s, Inches(2.4), Inches(1.95), Inches(8.5), Inches(0.9),
          "Skor = min(100,  B_max + 3 × n)", ukuran=32, warna=HIJAU,
          bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    _teks(s, Inches(2.4), Inches(3.15), Inches(8.5), Inches(0.6),
          "B_max = skor dasar severity tertinggi   •   n = jumlah temuan",
          ukuran=16, warna=RGBColor(0xCF, 0xDD, 0xF0), align=PP_ALIGN.CENTER)
    sev = [("AMAN", "0", HIJAU), ("RENDAH", "25", RGBColor(0x7C, 0xB3, 0x42)),
           ("SEDANG", "50", ORANYE), ("TINGGI", "75", RGBColor(0xE0, 0x6A, 0x1E)),
           ("KRITIS", "90", MERAH)]
    for i, (nama, skor, warna) in enumerate(sev):
        x = Inches(0.8 + i * 2.45)
        _kotak(s, x, Inches(4.5), Inches(2.2), Inches(1.7), warna)
        _teks(s, x, Inches(4.75), Inches(2.2), Inches(0.6), nama, ukuran=20,
              warna=PUTIH, bold=True, align=PP_ALIGN.CENTER)
        _teks(s, x, Inches(5.35), Inches(2.2), Inches(0.7), skor, ukuran=30,
              warna=PUTIH, bold=True, align=PP_ALIGN.CENTER)
    _teks(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6),
          "Banner usang menambah kategori I5 dan menaikkan severity satu "
          "tingkat.", ukuran=15, warna=RGBColor(0x9F, 0xB3, 0xD1),
          italic=True, align=PP_ALIGN.CENTER)


def slide_evolusi(prs):
    s = _slide_kosong(prs)
    _set_bg(s, PUTIH)
    _header(s, "Evolusi Prototipe: Awal → Saat Ini", 8)
    # Dua kolom perbandingan
    _kotak(s, Inches(0.8), Inches(1.45), Inches(5.6), Inches(0.7), ABU_TERANG)
    _teks(s, Inches(0.8), Inches(1.45), Inches(5.6), Inches(0.7),
          "VERSI AWAL (Prototipe M1)", ukuran=18, warna=PUTIH, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _kotak(s, Inches(6.9), Inches(1.45), Inches(5.6), Inches(0.7), HIJAU)
    _teks(s, Inches(6.9), Inches(1.45), Inches(5.6), Inches(0.7),
          "SAAT INI (M12)", ukuran=18, warna=PUTIH, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    baris = [
        ("TCP connect saja", "TCP + UDP (SSDP/UPnP, mDNS)"),
        ("Risiko dari port saja", "Port + kredensial + TLS + header HTTP"),
        ("IP + MAC", "+ vendor OUI + tipe perangkat"),
        ("OWASP I1/I2/I5", "OWASP I1–I10 lengkap"),
        ("Laporan JSON", "JSON + PDF + ringkasan OWASP"),
        ("CLI dasar", "CLI berwarna + web Streamlit"),
        ("~600 baris • 12 tes", "~2.287 baris • 36 tes + CI"),
    ]
    for i, (kiri, kanan) in enumerate(baris):
        y = Inches(2.3 + i * 0.66)
        bg = RGBColor(0xF2, 0xF5, 0xFA) if i % 2 == 0 else PUTIH
        _kotak(s, Inches(0.8), y, Inches(5.6), Inches(0.6), bg)
        _teks(s, Inches(1.0), y, Inches(5.2), Inches(0.6), kiri, ukuran=15,
              warna=ABU_TERANG, anchor=MSO_ANCHOR.MIDDLE)
        _kotak(s, Inches(6.9), y, Inches(5.6), Inches(0.6),
               RGBColor(0xE8, 0xF5, 0xEC) if i % 2 == 0 else PUTIH)
        _teks(s, Inches(7.1), y, Inches(5.2), Inches(0.6), kanan, ukuran=15,
              warna=ABU, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _teks(s, Inches(0.8), Inches(6.95), Inches(11.7), Inches(0.4),
          "Filosofi tetap: tanpa root, tanpa dependensi wajib, hanya jaringan "
          "sendiri — kedalaman analisis yang bertambah.",
          ukuran=14, warna=ABU_TERANG, italic=True, align=PP_ALIGN.CENTER)


def slide_contoh(prs):
    s = _slide_kosong(prs)
    _set_bg(s, PUTIH)
    _header(s, "Contoh Nyata: Kamera IP dengan Telnet", 9)
    _kotak(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.3),
           RGBColor(0xF2, 0xF5, 0xFA), garis=ABU_TERANG)
    _teks(s, Inches(1.0), Inches(1.7), Inches(5.2), Inches(0.5),
          "M1 (Awal)", ukuran=20, warna=ABU_TERANG, bold=True)
    _teks(s, Inches(1.0), Inches(2.3), Inches(5.2), Inches(1.4),
          "Terdeteksi port 2323 terbuka → kategori I1, severity TINGGI "
          "(hanya dari port).", ukuran=16, warna=ABU)
    _kotak(s, Inches(6.9), Inches(1.5), Inches(5.6), Inches(4.6),
           BIRU_TUA, garis=MERAH)
    _teks(s, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.5),
          "M12 (Saat Ini)", ukuran=20, warna=HIJAU, bold=True)
    _poin(s, Inches(7.1), Inches(2.35), Inches(5.2), Inches(3.6), [
        "Port 2323 terbuka (Telnet — vektor Mirai)",
        "Kredensial admin/admin BERHASIL (credcheck)",
        "Dikenali: Kamera IP (vendor Hikvision)",
        "HTTP tanpa header HSTS (tlscheck)",
    ], ukuran=16, warna=PUTIH, gap=10, warna_bullet=MERAH)
    _kotak(s, Inches(0.8), Inches(4.3), Inches(5.6), Inches(1.8), MERAH)
    _teks(s, Inches(0.8), Inches(4.55), Inches(5.6), Inches(0.6),
          "Hasil M12", ukuran=18, warna=PUTIH, bold=True,
          align=PP_ALIGN.CENTER)
    _teks(s, Inches(0.8), Inches(5.1), Inches(5.6), Inches(0.9),
          "KRITIS — skor 100 • I1/I2/I3/I9", ukuran=22, warna=PUTIH,
          bold=True, align=PP_ALIGN.CENTER)
    _teks(s, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.6),
          "Bukti empiris vektor serangan Mirai — persis sesuai teori pada "
          "Bab II paper.", ukuran=15, warna=ABU_TERANG, italic=True,
          align=PP_ALIGN.CENTER)


def slide_justifikasi(prs):
    s = _slide_kosong(prs)
    _set_bg(s, PUTIH)
    _header(s, "Justifikasi: Mengapa Aplikasi Ini Perlu Ada?", 10)
    _teks(s, Inches(0.8), Inches(1.25), Inches(11.7), Inches(0.6),
          "\"Bukankah sudah banyak alat serupa?\" — Kontribusi bukan teknik "
          "scanning baru, melainkan INTEGRASI unik:", ukuran=17, warna=ABU,
          bold=True)
    # Tabel pembanding sederhana
    kolom = ["Alat", "Open-\nsource", "Lokal\n(no cloud)", "Pemetaan\nOWASP IoT",
             "Ramah\nawam", "Bhs.\nIndonesia"]
    baris = [
        ("Nmap", "✓", "✓", "✗", "✗", "✗"),
        ("Nessus / OpenVAS", "~", "✓", "✗", "✗", "✗"),
        ("Fing", "✗", "✗", "✗", "✓", "✗"),
        ("Bitdefender", "✗", "✗", "✗", "✓", "✗"),
        ("HomeGuard", "✓", "✓", "✓", "✓", "✓"),
    ]
    x0, y0 = Inches(0.8), Inches(2.1)
    wcol = [Inches(3.0)] + [Inches(1.74)] * 5
    hhead = Inches(0.75)
    hrow = Inches(0.72)
    # Header
    cx = x0
    for j, k in enumerate(kolom):
        _kotak(s, cx, y0, wcol[j], hhead, BIRU_TUA)
        _teks(s, cx, y0, wcol[j], hhead, k, ukuran=13, warna=PUTIH, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx = Inches(cx.inches + wcol[j].inches)
    # Baris
    for i, row in enumerate(baris):
        ry = Inches(y0.inches + hhead.inches + i * hrow.inches)
        is_hg = row[0] == "HomeGuard"
        cx = x0
        for j, val in enumerate(row):
            if j == 0:
                bg = HIJAU if is_hg else RGBColor(0xEC, 0xEF, 0xF4)
                fg = PUTIH if is_hg else ABU
                _kotak(s, cx, ry, wcol[j], hrow, bg)
                _teks(s, cx, ry, wcol[j], hrow, val, ukuran=14, warna=fg,
                      bold=is_hg, align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)
            else:
                if val == "✓":
                    bg, fg = (RGBColor(0xD6, 0xF0, 0xDE), HIJAU)
                elif val == "✗":
                    bg, fg = (RGBColor(0xFB, 0xE2, 0xE2), MERAH)
                else:
                    bg, fg = (RGBColor(0xFC, 0xF1, 0xDD), ORANYE)
                if is_hg:
                    bg = RGBColor(0xE8, 0xF5, 0xEC)
                _kotak(s, cx, ry, wcol[j], hrow, bg)
                _teks(s, cx, ry, wcol[j], hrow, val, ukuran=18, warna=fg,
                      bold=True, align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)
            cx = Inches(cx.inches + wcol[j].inches)
    _teks(s, Inches(0.8), Inches(6.75), Inches(11.7), Inches(0.5),
          "HomeGuard = satu-satunya yang memenuhi kelima kriteria sekaligus.",
          ukuran=16, warna=HIJAU, bold=True, align=PP_ALIGN.CENTER)


def slide_pengujian(prs):
    s = _slide_kosong(prs)
    _set_bg(s, PUTIH)
    _header(s, "Pengujian & Jaminan Kualitas", 11)
    stat = [
        ("36", "Unit test\n(semua lulus)", AKSEN),
        ("4", "Versi Python\n(3.9 – 3.12)", HIJAU),
        ("12", "Modul Python\nteruji", ORANYE),
        ("~2.287", "Baris kode\ninti", BIRU),
    ]
    for i, (angka, label, warna) in enumerate(stat):
        x = Inches(0.8 + i * 3.05)
        _kotak(s, x, Inches(1.6), Inches(2.8), Inches(2.0), warna)
        _teks(s, x, Inches(1.85), Inches(2.8), Inches(0.9), angka, ukuran=44,
              warna=PUTIH, bold=True, align=PP_ALIGN.CENTER)
        _teks(s, x, Inches(2.85), Inches(2.8), Inches(0.7), label, ukuran=15,
              warna=PUTIH, align=PP_ALIGN.CENTER)
    _teks(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.5),
          "Cakupan pengujian per modul:", ukuran=18, warna=ABU, bold=True)
    _poin(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), [
        "discovery, portscan, vulnmap, scanner — alur pipeline inti",
        "classify (tipe perangkat) & tlscheck (TLS/cipher/header) — modul "
        "lanjutan",
        "GitHub Actions CI menjalankan seluruh tes otomatis tiap perubahan "
        "kode",
        "Demo aman: python demo_lokal.py — tanpa jaringan nyata",
    ], ukuran=17, gap=12)


def slide_demo(prs):
    s = _slide_kosong(prs)
    _set_bg(s, BIRU_TUA)
    _kotak(s, 0, 0, LEBAR, Inches(1.05), BIRU)
    _kotak(s, 0, Inches(1.05), LEBAR, Inches(0.06), HIJAU)
    _teks(s, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
          "Demo Langsung", ukuran=28, warna=PUTIH, bold=True,
          anchor=MSO_ANCHOR.MIDDLE)
    _teks(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5),
          "Tiga cara membuktikan aplikasi bekerja:", ukuran=18,
          warna=RGBColor(0xCF, 0xDD, 0xF0))
    cmds = [
        ("Demo aman (tanpa jaringan nyata)", "python demo_lokal.py"),
        ("Antarmuka web interaktif", "streamlit run app.py"),
        ("Pindai dari terminal", "python cli.py --subnet 192.168.1.0/24 "
         "--check-creds --tls --pdf"),
        ("Lihat referensi OWASP IoT", "python cli.py --owasp-ref"),
    ]
    for i, (label, cmd) in enumerate(cmds):
        y = Inches(2.3 + i * 1.15)
        _teks(s, Inches(0.8), y, Inches(11.7), Inches(0.4), label, ukuran=16,
              warna=HIJAU, bold=True)
        _kotak(s, Inches(0.8), Inches(y.inches + 0.42), Inches(11.7),
               Inches(0.6), RGBColor(0x05, 0x10, 0x24))
        _teks(s, Inches(1.1), Inches(y.inches + 0.42), Inches(11.2),
              Inches(0.6), "$ " + cmd, ukuran=17,
              warna=RGBColor(0x6F, 0xE3, 0x9F), font="Consolas",
              anchor=MSO_ANCHOR.MIDDLE)


def slide_kesimpulan(prs):
    s = _slide_kosong(prs)
    _set_bg(s, PUTIH)
    _header(s, "Kesimpulan & Roadmap", 13)
    _teks(s, Inches(0.8), Inches(1.35), Inches(5.7), Inches(0.5),
          "Kesimpulan", ukuran=22, warna=HIJAU, bold=True)
    _poin(s, Inches(0.8), Inches(2.0), Inches(5.7), Inches(4.5), [
        "Prototipe HomeGuard berhasil dibangun sesuai rancangan paper "
        "(pipeline 4 modul)",
        "Memetakan kerentanan IoT ke OWASP IoT Top 10 dengan skor risiko "
        "yang mudah dipahami",
        "Berjalan lokal, tanpa root, tanpa dependensi wajib — ramah pengguna "
        "awam",
        "Terbukti mendeteksi vektor serangan nyata (mis. Mirai) secara "
        "empiris",
    ], ukuran=17, gap=14, warna_bullet=HIJAU)
    _teks(s, Inches(6.9), Inches(1.35), Inches(5.7), Inches(0.5),
          "Roadmap", ukuran=22, warna=ORANYE, bold=True)
    _poin(s, Inches(6.9), Inches(2.0), Inches(5.7), Inches(4.5), [
        "Korelasi versi komponen → CVE nyata via NVD API",
        "Pengujian dunia nyata pada jaringan rumah + isi tabel hasil paper",
        "Screenshot antarmuka Streamlit untuk lampiran",
        "Dukungan IPv6 & packaging (pyproject.toml)",
    ], ukuran=17, gap=14, warna_bullet=ORANYE)


def slide_penutup(prs):
    s = _slide_kosong(prs)
    _set_bg(s, BIRU_TUA)
    _kotak(s, 0, Inches(3.0), LEBAR, Inches(0.08), HIJAU)
    _teks(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(0.8),
          "Terima Kasih", ukuran=48, warna=PUTIH, bold=True,
          align=PP_ALIGN.CENTER)
    _teks(s, Inches(1.0), Inches(3.3), Inches(11.3), Inches(0.6),
          "Diskusi & Tanya Jawab", ukuran=22, warna=HIJAU,
          align=PP_ALIGN.CENTER, italic=True)
    _teks(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.5),
          "Reiza Gerrard Rizki Ramadhan  •  Politeknik Siber dan Sandi Negara",
          ukuran=16, warna=RGBColor(0x9F, 0xB3, 0xD1), align=PP_ALIGN.CENTER)
    _teks(s, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.5),
          "HomeGuard — Pemindai Keamanan IoT Berbasis OWASP IoT Top 10",
          ukuran=14, warna=RGBColor(0x6B, 0x7A, 0x95), align=PP_ALIGN.CENTER,
          italic=True)


def main() -> int:
    prs = Presentation()
    prs.slide_width = LEBAR
    prs.slide_height = TINGGI

    slide_judul(prs)
    slide_agenda(prs)
    slide_latar(prs)
    slide_tujuan(prs)
    slide_arsitektur(prs)
    slide_caramodul(prs)
    slide_rumus(prs)
    slide_evolusi(prs)
    slide_contoh(prs)
    slide_justifikasi(prs)
    slide_pengujian(prs)
    slide_demo(prs)
    slide_kesimpulan(prs)
    slide_penutup(prs)

    prs.save(TUJUAN)
    print(f"Berhasil: {TUJUAN} ({os.path.getsize(TUJUAN)} byte, "
          f"{len(prs.slides._sldIdLst)} slide)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
